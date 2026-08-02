# -*- coding: utf-8 -*-
# text_studio_app/lib/slide_creation/layouts/section_slides.py
# ============================================================
# 見出しページ描画
# ============================================================

from __future__ import annotations

from typing import Any

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

from lib.slide_creation.components import (
    add_textbox,
    set_slide_background,
)
from lib.slide_creation.models import SlideDefinition, SlideTheme
from lib.slide_creation.theme_image_renderer import (
    add_theme_background_image,
)
from lib.slide_creation.theme_layouts.registry import (
    render_registered_section,
)

def add_section_slide(
    prs: Any,
    *,
    slide_def: SlideDefinition,
    theme: SlideTheme,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_theme_background_image(
        slide=slide,
        prs=prs,
        theme=theme,
        slide_type="section",
    )

    style_key = slide_def.style_key

    # ========================================================
    # 専用レイアウト
    # ========================================================
    if render_registered_section(
        slide,
        slide_def=slide_def,
        theme=theme,
    ):
        return
    
    if style_key == "full_color":
        set_slide_background(slide, color=theme.primary_color)

        add_textbox(
            slide,
            left=1.0,
            top=1.6,
            width=11.3,
            height=0.8,
            text=slide_def.section_number,
            theme=theme,
            font_size=26,
            color=theme.accent_color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        add_textbox(
            slide,
            left=1.0,
            top=2.5,
            width=11.3,
            height=1.4,
            text=slide_def.title,
            theme=theme,
            font_size=40,
            color=theme.title_text_color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

        add_textbox(
            slide,
            left=1.5,
            top=4.05,
            width=10.3,
            height=0.8,
            text=slide_def.subtitle,
            theme=theme,
            font_size=17,
            color=theme.title_text_color,
            alignment=PP_ALIGN.CENTER,
        )
        return

    set_slide_background(slide, color=theme.background_color)

    if style_key == "large_number":
        add_textbox(
            slide,
            left=0.8,
            top=1.0,
            width=4.0,
            height=2.5,
            text=slide_def.section_number,
            theme=theme,
            font_size=80,
            color=theme.accent_color,
            bold=True,
        )

        add_textbox(
            slide,
            left=3.6,
            top=2.2,
            width=8.6,
            height=1.2,
            text=slide_def.title,
            theme=theme,
            font_size=38,
            color=theme.body_text_color,
            bold=True,
        )

        add_textbox(
            slide,
            left=3.65,
            top=3.55,
            width=8.2,
            height=0.8,
            text=slide_def.subtitle,
            theme=theme,
            font_size=17,
            color=theme.sub_text_color,
        )
        return

    if style_key == "centered":
        add_textbox(
            slide,
            left=1.0,
            top=1.8,
            width=11.3,
            height=0.6,
            text=slide_def.section_number,
            theme=theme,
            font_size=22,
            color=theme.accent_color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        add_textbox(
            slide,
            left=1.0,
            top=2.55,
            width=11.3,
            height=1.2,
            text=slide_def.title,
            theme=theme,
            font_size=40,
            color=theme.body_text_color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        add_textbox(
            slide,
            left=1.4,
            top=3.9,
            width=10.5,
            height=0.8,
            text=slide_def.subtitle,
            theme=theme,
            font_size=17,
            color=theme.sub_text_color,
            alignment=PP_ALIGN.CENTER,
        )
        return

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.85),
        Inches(1.25),
        Inches(0.12),
        Inches(4.8),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = theme.accent_color
    accent.line.fill.background()

    add_textbox(
        slide,
        left=1.3,
        top=1.6,
        width=10.5,
        height=0.6,
        text=slide_def.section_number,
        theme=theme,
        font_size=22,
        color=theme.accent_color,
        bold=True,
    )

    add_textbox(
        slide,
        left=1.3,
        top=2.35,
        width=10.5,
        height=1.2,
        text=slide_def.title,
        theme=theme,
        font_size=38,
        color=theme.body_text_color,
        bold=True,
    )

    add_textbox(
        slide,
        left=1.3,
        top=3.75,
        width=10.3,
        height=0.9,
        text=slide_def.subtitle,
        theme=theme,
        font_size=17,
        color=theme.sub_text_color,
    )