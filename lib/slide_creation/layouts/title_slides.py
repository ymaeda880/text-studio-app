# -*- coding: utf-8 -*-
# text_studio_app/lib/slide_creation/layouts/title_slides.py
# ============================================================
# タイトルページ描画
# ============================================================

from __future__ import annotations

from typing import Any

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

from lib.slide_creation.components import (
    add_textbox,
    set_slide_background,
)
from lib.slide_creation.models import (
    PresentationSettings,
    SlideDefinition,
    SlideTheme,
)
from lib.slide_creation.theme_image_renderer import (
    add_theme_background_image,
)
from lib.slide_creation.theme_layouts.registry import (
    render_registered_title,
)

def add_title_slide(
    prs: Any,
    *,
    slide_def: SlideDefinition,
    settings: PresentationSettings,
    theme: SlideTheme,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_theme_background_image(
        slide=slide,
        prs=prs,
        theme=theme,
        slide_type="title",
    )

    style_key = slide_def.style_key

    if render_registered_title(
        slide,
        slide_def=slide_def,
        settings=settings,
        theme=theme,
    ):
        return

    if style_key == "full_color":
        set_slide_background(slide, color=theme.primary_color)

        add_textbox(
            slide,
            left=1.0,
            top=2.1,
            width=11.3,
            height=1.4,
            text=slide_def.title,
            theme=theme,
            font_size=38,
            color=theme.title_text_color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

        add_textbox(
            slide,
            left=1.3,
            top=3.6,
            width=10.7,
            height=0.7,
            text=slide_def.subtitle,
            theme=theme,
            font_size=20,
            color=theme.title_text_color,
            alignment=PP_ALIGN.CENTER,
        )

    elif style_key == "centered":
        set_slide_background(slide, color=theme.background_color)

        add_textbox(
            slide,
            left=1.0,
            top=2.15,
            width=11.3,
            height=1.4,
            text=slide_def.title,
            theme=theme,
            font_size=38,
            color=theme.body_text_color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

        add_textbox(
            slide,
            left=1.4,
            top=3.65,
            width=10.5,
            height=0.7,
            text=slide_def.subtitle,
            theme=theme,
            font_size=19,
            color=theme.sub_text_color,
            alignment=PP_ALIGN.CENTER,
        )

    elif style_key == "minimal":
        set_slide_background(slide, color=theme.background_color)

        add_textbox(
            slide,
            left=1.0,
            top=2.35,
            width=11.0,
            height=1.3,
            text=slide_def.title,
            theme=theme,
            font_size=36,
            color=theme.body_text_color,
            bold=True,
        )

        add_textbox(
            slide,
            left=1.0,
            top=3.75,
            width=11.0,
            height=0.7,
            text=slide_def.subtitle,
            theme=theme,
            font_size=18,
            color=theme.sub_text_color,
        )

    else:
        set_slide_background(slide, color=theme.background_color)

        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0,
            0,
            Inches(0.3),
            prs.slide_height,
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = theme.accent_color
        accent.line.fill.background()

        add_textbox(
            slide,
            left=1.0,
            top=2.15,
            width=11.2,
            height=1.35,
            text=slide_def.title,
            theme=theme,
            font_size=36,
            color=theme.body_text_color,
            bold=True,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

        add_textbox(
            slide,
            left=1.0,
            top=3.75,
            width=11.0,
            height=0.7,
            text=slide_def.subtitle,
            theme=theme,
            font_size=18,
            color=theme.sub_text_color,
        )

    detail_values = [
        value
        for value in (
            settings.company_name,
            slide_def.presenter_name or settings.presenter_name,
            settings.presentation_date,
        )
        if value
    ]

    add_textbox(
        slide,
        left=1.0,
        top=6.25,
        width=11.0,
        height=0.4,
        text="　｜　".join(detail_values),
        theme=theme,
        font_size=11,
        color=(
            theme.title_text_color
            if style_key == "full_color"
            else theme.sub_text_color
        ),
        alignment=PP_ALIGN.RIGHT,
    )