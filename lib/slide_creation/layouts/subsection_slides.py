# -*- coding: utf-8 -*-
# text_studio_app/lib/slide_creation/layouts/subsection_slides.py
# ============================================================
# 小見出しページ描画
# ============================================================

from __future__ import annotations

from typing import Any

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

from lib.slide_creation.components import (
    add_panel,
    add_textbox,
    set_slide_background,
)
from lib.slide_creation.models import SlideDefinition, SlideTheme
from lib.slide_creation.theme_image_renderer import (
    add_theme_background_image,
)
from lib.slide_creation.theme_layouts.registry import (
    render_registered_subsection,
)

def add_subsection_slide(
    prs: Any,
    *,
    slide_def: SlideDefinition,
    theme: SlideTheme,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_theme_background_image(
        slide=slide,
        prs=prs,
        theme=theme,
        slide_type="subsection",
    )

    style_key = slide_def.style_key

    if render_registered_subsection(
        slide,
        slide_def=slide_def,
        theme=theme,
    ):
        return

    set_slide_background(
        slide,
        color=theme.background_color,
    )


    if style_key == "accent_band":
        band = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0,
            Inches(2.2),
            prs.slide_width,
            Inches(2.1),
        )
        band.fill.solid()
        band.fill.fore_color.rgb = theme.primary_color
        band.line.fill.background()

        add_textbox(
            slide,
            left=1.0,
            top=2.45,
            width=11.3,
            height=0.9,
            text=slide_def.title,
            theme=theme,
            font_size=34,
            color=theme.title_text_color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

        add_textbox(
            slide,
            left=1.5,
            top=3.35,
            width=10.3,
            height=0.6,
            text=slide_def.subtitle,
            theme=theme,
            font_size=16,
            color=theme.title_text_color,
            alignment=PP_ALIGN.CENTER,
        )
        return

    if style_key == "card":
        add_panel(
            slide,
            left=1.15,
            top=1.65,
            width=11.0,
            height=4.1,
            theme=theme,
        )

        add_textbox(
            slide,
            left=1.65,
            top=2.25,
            width=10.0,
            height=1.0,
            text=slide_def.title,
            theme=theme,
            font_size=34,
            color=theme.body_text_color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        add_textbox(
            slide,
            left=2.0,
            top=3.55,
            width=9.3,
            height=0.8,
            text=slide_def.subtitle,
            theme=theme,
            font_size=17,
            color=theme.sub_text_color,
            alignment=PP_ALIGN.CENTER,
        )
        return

    if style_key == "dark_panel":
        add_panel(
            slide,
            left=0.85,
            top=1.55,
            width=11.65,
            height=4.3,
            theme=theme,
            fill_color=theme.primary_color,
        )

        add_textbox(
            slide,
            left=1.35,
            top=2.25,
            width=10.7,
            height=1.0,
            text=slide_def.title,
            theme=theme,
            font_size=34,
            color=theme.title_text_color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        add_textbox(
            slide,
            left=1.7,
            top=3.55,
            width=10.0,
            height=0.8,
            text=slide_def.subtitle,
            theme=theme,
            font_size=17,
            color=theme.title_text_color,
            alignment=PP_ALIGN.CENTER,
        )
        return

    add_textbox(
        slide,
        left=1.0,
        top=2.4,
        width=11.3,
        height=1.0,
        text=slide_def.title,
        theme=theme,
        font_size=36,
        color=theme.body_text_color,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    add_textbox(
        slide,
        left=1.5,
        top=3.65,
        width=10.3,
        height=0.8,
        text=slide_def.subtitle,
        theme=theme,
        font_size=17,
        color=theme.sub_text_color,
        alignment=PP_ALIGN.CENTER,
    )