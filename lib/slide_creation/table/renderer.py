# -*- coding: utf-8 -*-
# text_studio_app/lib/slide_creation/table/renderer.py
# ============================================================
# SlideTex 表描画
#
# 機能：
# - TableDefinitionからPowerPoint表を作成する
# - 表題を表示する
# - ヘッダー行・ヘッダー列へスタイルを適用する
# - セル内改行を表示する
# - <同上>・<同左>によるセル結合を行う
# - 表スタイルに応じて背景色・文字・罫線を設定する
#
# 方針：
# - SlideTex構文解析には依存しない
# - 表データはTableDefinitionから受け取る
# - 色とフォントはSlideThemeを使用する
# - 表の配置領域は呼出元から受け取る
# ============================================================

from __future__ import annotations

# ============================================================
# imports
# ============================================================
from collections import defaultdict
from typing import Any

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

from lib.slide_creation.font_settings import resolve_font_size
from lib.slide_creation.models import (
    SlideTheme,
)
from lib.slide_creation.table.models import (
    TableCell,
    TableDefinition,
)
from lib.slide_creation.table.styles import (
    TableStyleDefinition,
    get_table_style,
)


# ============================================================
# 定数
# ============================================================
DEFAULT_CAPTION_HEIGHT = 0.38
MIN_TABLE_HEIGHT = 0.6


# ============================================================
# 色
# ============================================================
def _rgb_to_hex(
    color: Any,
) -> str:
    # ------------------------------------------------------------
    # RGBColor等を6桁の16進文字列へ変換する
    # ------------------------------------------------------------
    if color is None:
        return "000000"

    try:
        return (
            f"{int(color[0]):02X}"
            f"{int(color[1]):02X}"
            f"{int(color[2]):02X}"
        )

    except (
        TypeError,
        ValueError,
        IndexError,
    ):
        return "000000"


# ============================================================
# セル取得
# ============================================================
def _get_definition_cell(
    table_definition: TableDefinition,
    *,
    row_index: int,
    column_index: int,
) -> TableCell | None:
    if row_index < 0 or column_index < 0:
        return None

    if row_index >= len(table_definition.rows):
        return None

    row = table_definition.rows[row_index]

    if column_index >= len(row):
        return None

    return row[column_index]


# ============================================================
# セル結合元
# ============================================================
def _resolve_origin_position(
    table_definition: TableDefinition,
    *,
    row_index: int,
    column_index: int,
) -> tuple[int, int] | None:
    # ------------------------------------------------------------
    # <同上>・<同左>が最終的に参照する通常セルを取得する
    # ------------------------------------------------------------
    current_row = row_index
    current_column = column_index

    visited: set[tuple[int, int]] = set()

    while True:
        position = (
            current_row,
            current_column,
        )

        if position in visited:
            return None

        visited.add(
            position,
        )

        cell = _get_definition_cell(
            table_definition,
            row_index=current_row,
            column_index=current_column,
        )

        if cell is None:
            return None

        if cell.merge_direction == "":
            return (
                current_row,
                current_column,
            )

        if cell.merge_direction == "up":
            current_row -= 1
            continue

        if cell.merge_direction == "left":
            current_column -= 1
            continue

        return None


def _build_merge_ranges(
    table_definition: TableDefinition,
) -> list[tuple[int, int, int, int]]:
    # ------------------------------------------------------------
    # 同じ基準セルへ属するセルをまとめ，
    # PowerPointで結合する矩形範囲へ変換する
    # ------------------------------------------------------------
    groups: dict[
        tuple[int, int],
        set[tuple[int, int]],
    ] = defaultdict(set)

    for row_index, row in enumerate(
        table_definition.rows,
    ):
        for column_index, _cell in enumerate(
            row,
        ):
            origin = _resolve_origin_position(
                table_definition,
                row_index=row_index,
                column_index=column_index,
            )

            if origin is None:
                continue

            groups[origin].add(
                (
                    row_index,
                    column_index,
                )
            )

    merge_ranges: list[
        tuple[int, int, int, int]
    ] = []

    for origin, positions in groups.items():
        if len(positions) <= 1:
            continue

        row_indexes = [
            position[0]
            for position in positions
        ]

        column_indexes = [
            position[1]
            for position in positions
        ]

        start_row = min(
            row_indexes,
        )

        end_row = max(
            row_indexes,
        )

        start_column = min(
            column_indexes,
        )

        end_column = max(
            column_indexes,
        )

        merge_ranges.append(
            (
                start_row,
                start_column,
                end_row,
                end_column,
            )
        )

    return merge_ranges


# ============================================================
# ヘッダー判定
# ============================================================
def _is_header_cell(
    table_definition: TableDefinition,
    *,
    row_index: int,
    column_index: int,
) -> bool:
    return (
        row_index
        < table_definition.header.row_count
        or column_index
        < table_definition.header.column_count
    )


def _is_top_header_cell(
    table_definition: TableDefinition,
    *,
    row_index: int,
) -> bool:
    return (
        row_index
        < table_definition.header.row_count
    )


def _is_row_header_cell(
    table_definition: TableDefinition,
    *,
    column_index: int,
) -> bool:
    return (
        column_index
        < table_definition.header.column_count
    )


# ============================================================
# セル背景
# ============================================================
def _set_cell_fill(
    cell: Any,
    color: Any,
) -> None:
    if color is None:
        cell.fill.background()
        return

    cell.fill.solid()
    cell.fill.fore_color.rgb = color


# ============================================================
# セル罫線
# ============================================================
def _remove_xml_children(
    element: Any,
    tag_name: str,
) -> None:
    # ------------------------------------------------------------
    # 指定されたローカル名のXML要素を削除する
    #
    # child.tag は
    # {http://schemas.openxmlformats.org/drawingml/2006/main}lnL
    # の形式になるため，a:lnL のような接頭辞を除いて比較する
    # ------------------------------------------------------------
    local_name = (
        tag_name.split(
            ":",
            maxsplit=1,
        )[-1]
    )

    for child in list(
        element,
    ):
        child_local_name = (
            str(
                child.tag,
            ).split(
                "}",
                maxsplit=1,
            )[-1]
        )

        if child_local_name == local_name:
            element.remove(
                child,
            )


def _set_cell_border_side(
    cell: Any,
    *,
    side: str,
    color: Any,
    width_pt: float,
    visible: bool,
) -> None:
    # ------------------------------------------------------------
    # python-pptxにはセル罫線の公開APIがないため，
    # DrawingMLのXMLを直接操作する
    #
    # side:
    # - a:lnL：左罫線
    # - a:lnR：右罫線
    # - a:lnT：上罫線
    # - a:lnB：下罫線
    # ------------------------------------------------------------
    tc = cell._tc
    tc_properties = tc.get_or_add_tcPr()

    # --------------------------------------------------------
    # 同じ辺の既存罫線を削除する
    #
    # 同じlnL等が複数存在すると，PowerPointによって
    # 罫線が正しく解釈されない場合がある
    # --------------------------------------------------------
    _remove_xml_children(
        tc_properties,
        side,
    )

    # --------------------------------------------------------
    # 罫線要素を作成する
    #
    # DrawingMLの線幅はEMU単位
    # 1pt = 12700 EMU
    # --------------------------------------------------------
    line = OxmlElement(
        side,
    )

    line.set(
        "w",
        str(
            int(
                max(
                    float(
                        width_pt,
                    ),
                    0.1,
                )
                * 12700
            )
        ),
    )

    line.set(
        "cap",
        "flat",
    )

    line.set(
        "cmpd",
        "sng",
    )

    line.set(
        "algn",
        "ctr",
    )

    if visible:
        # ----------------------------------------------------
        # 線色
        # ----------------------------------------------------
        solid_fill = OxmlElement(
            "a:solidFill",
        )

        srgb_color = OxmlElement(
            "a:srgbClr",
        )

        srgb_color.set(
            "val",
            _rgb_to_hex(
                color,
            ),
        )

        solid_fill.append(
            srgb_color,
        )

        line.append(
            solid_fill,
        )

        # ----------------------------------------------------
        # 実線
        # ----------------------------------------------------
        dash = OxmlElement(
            "a:prstDash",
        )

        dash.set(
            "val",
            "solid",
        )

        line.append(
            dash,
        )

        # ----------------------------------------------------
        # 線の結合形状
        # ----------------------------------------------------
        round_join = OxmlElement(
            "a:round",
        )

        line.append(
            round_join,
        )

        # ----------------------------------------------------
        # 線端
        # ----------------------------------------------------
        head_end = OxmlElement(
            "a:headEnd",
        )

        head_end.set(
            "type",
            "none",
        )

        head_end.set(
            "w",
            "med",
        )

        head_end.set(
            "len",
            "med",
        )

        line.append(
            head_end,
        )

        tail_end = OxmlElement(
            "a:tailEnd",
        )

        tail_end.set(
            "type",
            "none",
        )

        tail_end.set(
            "w",
            "med",
        )

        tail_end.set(
            "len",
            "med",
        )

        line.append(
            tail_end,
        )

    else:
        # ----------------------------------------------------
        # 非表示の辺
        # ----------------------------------------------------
        no_fill = OxmlElement(
            "a:noFill",
        )

        line.append(
            no_fill,
        )

    tc_properties.append(
        line,
    )

def _set_cell_borders(
    cell: Any,
    *,
    row_index: int,
    column_index: int,
    row_count: int,
    column_count: int,
    style: TableStyleDefinition,
    border_color: Any,
) -> None:
    is_first_row = row_index == 0
    is_last_row = row_index == row_count - 1
    is_first_column = column_index == 0
    is_last_column = column_index == column_count - 1

    left_visible = (
        style.show_outer_border
        if is_first_column
        else style.show_inner_vertical_border
    )

    right_visible = (
        style.show_outer_border
        if is_last_column
        else style.show_inner_vertical_border
    )

    top_visible = (
        style.show_outer_border
        if is_first_row
        else style.show_inner_horizontal_border
    )

    bottom_visible = (
        style.show_outer_border
        if is_last_row
        else style.show_inner_horizontal_border
    )

    _set_cell_border_side(
        cell,
        side="a:lnL",
        color=border_color,
        width_pt=style.border_width_pt,
        visible=left_visible,
    )

    _set_cell_border_side(
        cell,
        side="a:lnR",
        color=border_color,
        width_pt=style.border_width_pt,
        visible=right_visible,
    )

    _set_cell_border_side(
        cell,
        side="a:lnT",
        color=border_color,
        width_pt=style.border_width_pt,
        visible=top_visible,
    )

    _set_cell_border_side(
        cell,
        side="a:lnB",
        color=border_color,
        width_pt=style.border_width_pt,
        visible=bottom_visible,
    )


# ============================================================
# セル文字
# ============================================================
def _set_cell_text(
    cell: Any,
    *,
    text: str,
    theme: SlideTheme,
    font_size: int,
    font_name: str | None,
    color: Any,
    bold: bool,
    alignment: PP_ALIGN,
) -> None:
    frame = cell.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    lines = str(
        text or ""
    ).splitlines()

    if not lines:
        lines = [
            "",
        ]

    for line_index, line in enumerate(
        lines,
    ):
        paragraph = (
            frame.paragraphs[0]
            if line_index == 0
            else frame.add_paragraph()
        )

        paragraph.text = line
        paragraph.alignment = alignment
        paragraph.space_after = Pt(0)
        paragraph.space_before = Pt(0)

        for run in paragraph.runs:
            run.font.name = font_name or theme.font_name
            run.font.size = Pt(
                font_size,
            )
            run.font.bold = bold
            run.font.color.rgb = color


# ============================================================
# セル余白
# ============================================================
def _set_cell_margins(
    cell: Any,
    *,
    style: TableStyleDefinition,
) -> None:
    cell.margin_left = Inches(
        style.margin_left_inch,
    )

    cell.margin_right = Inches(
        style.margin_right_inch,
    )

    cell.margin_top = Inches(
        style.margin_top_inch,
    )

    cell.margin_bottom = Inches(
        style.margin_bottom_inch,
    )


# ============================================================
# セルスタイル
# ============================================================
def _get_cell_fill_color(
    table_definition: TableDefinition,
    *,
    row_index: int,
    column_index: int,
    style: TableStyleDefinition,
) -> Any:
    # ------------------------------------------------------------
    # 上部ヘッダー行を優先し，
    # 次に左端ヘッダー列を適用する
    # ------------------------------------------------------------
    if _is_top_header_cell(
        table_definition,
        row_index=row_index,
    ):
        return style.header_fill_color

    if _is_row_header_cell(
        table_definition,
        column_index=column_index,
    ):
        return style.row_header_fill_color

    if (
        style.use_banded_rows
        and row_index % 2 == 1
    ):
        return style.banded_fill_color

    return style.body_fill_color


def _get_cell_text_color(
    table_definition: TableDefinition,
    *,
    row_index: int,
    column_index: int,
    style: TableStyleDefinition,
) -> Any:
    if _is_top_header_cell(
        table_definition,
        row_index=row_index,
    ):
        return style.header_text_color

    if _is_row_header_cell(
        table_definition,
        column_index=column_index,
    ):
        return style.row_header_text_color

    return style.body_text_color


# ============================================================
# セル配置
# ============================================================
def _get_cell_alignment(
    table_definition: TableDefinition,
    *,
    row_index: int,
    column_index: int,
    text: str,
) -> PP_ALIGN:
    # ------------------------------------------------------------
    # 上部ヘッダー
    # ------------------------------------------------------------
    if _is_top_header_cell(
        table_definition,
        row_index=row_index,
    ):
        return PP_ALIGN.CENTER

    # ------------------------------------------------------------
    # 左端ヘッダー列
    # ------------------------------------------------------------
    if _is_row_header_cell(
        table_definition,
        column_index=column_index,
    ):
        return PP_ALIGN.LEFT

    # ------------------------------------------------------------
    # 数値は右寄せ
    #
    # カンマ，符号，小数点，%を除いて
    # 数字だけになる場合を数値と判定する
    # ------------------------------------------------------------
    normalized_text = (
        str(text or "")
        .strip()
        .replace(",", "")
        .replace("，", "")
        .replace("%", "")
        .replace("％", "")
        .replace(".", "")
        .replace("-", "")
        .replace("+", "")
        .replace("△", "")
    )

    if (
        normalized_text
        and normalized_text.isdigit()
    ):
        return PP_ALIGN.RIGHT

    return PP_ALIGN.LEFT


# ============================================================
# 表題
# ============================================================
def _add_caption(
    slide: Any,
    *,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    theme: SlideTheme,
    font_name: str | None = None,
    font_size: int | None = None,
) -> None:
    # ------------------------------------------------------------
    # 左側アクセント線
    # ------------------------------------------------------------
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(top + 0.07),
        Inches(0.06),
        Inches(max(height - 0.14, 0.12)),
    )

    accent.fill.solid()
    accent.fill.fore_color.rgb = theme.accent_color
    accent.line.fill.background()

    # ------------------------------------------------------------
    # 表題
    # ------------------------------------------------------------
    caption_box = slide.shapes.add_textbox(
        Inches(left + 0.15),
        Inches(top),
        Inches(width - 0.15),
        Inches(height),
    )

    frame = caption_box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.LEFT
    paragraph.space_after = Pt(0)

    for run in paragraph.runs:
        run.font.name = font_name or theme.font_name
        run.font.size = Pt(font_size or theme.table_caption_font_size)
        run.font.bold = True
        run.font.color.rgb = theme.body_text_color


# ============================================================
# 行高
# ============================================================
def _set_uniform_row_heights(
    table: Any,
    *,
    total_height: float,
) -> None:
    row_count = len(
        table.rows,
    )

    if row_count <= 0:
        return

    row_height = (
        float(total_height)
        / float(row_count)
    )

    # ===== DEBUG START =====
    print("----- TABLE ROW HEIGHT -----")
    print(f"{total_height=}")
    print(f"{row_count=}")
    print(f"{row_height=}")
    # ===== DEBUG END =====

    for row in table.rows:
        row.height = Inches(
            row_height,
        )


# ============================================================
# 列幅
# ============================================================
def _set_uniform_column_widths(
    table: Any,
    *,
    total_width: float,
    header_column_count: int = 0,
) -> None:
    column_count = len(
        table.columns,
    )

    if column_count <= 0:
        return

    # --------------------------------------------------------
    # 左端にヘッダー列がある場合
    # - 左端列を少し狭くする
    # - 残りを均等に配分する
    # --------------------------------------------------------
    if (
        header_column_count >= 1
        and column_count >= 2
    ):
        first_column_width = (
            float(total_width)
            * 0.22
        )

        remaining_width = (
            float(total_width)
            - first_column_width
        )

        other_column_width = (
            remaining_width
            / float(column_count - 1)
        )

        table.columns[0].width = Inches(
            first_column_width,
        )

        for column_index in range(
            1,
            column_count,
        ):
            table.columns[
                column_index
            ].width = Inches(
                other_column_width,
            )

        return

    # --------------------------------------------------------
    # ヘッダー列がない場合は均等幅
    # --------------------------------------------------------
    column_width = (
        float(total_width)
        / float(column_count)
    )

    for column in table.columns:
        column.width = Inches(
            column_width,
        )


# ============================================================
# 結合
# ============================================================
def _merge_table_cells(
    pptx_table: Any,
    *,
    table_definition: TableDefinition,
) -> None:
    merge_ranges = _build_merge_ranges(
        table_definition,
    )

    for (
        start_row,
        start_column,
        end_row,
        end_column,
    ) in merge_ranges:
        start_cell = pptx_table.cell(
            start_row,
            start_column,
        )

        end_cell = pptx_table.cell(
            end_row,
            end_column,
        )

        start_cell.merge(
            end_cell,
        )


# ============================================================
# 公開関数
# ============================================================
def render_table(
    slide: Any,
    *,
    table_definition: TableDefinition,
    theme: SlideTheme,
    left: float,
    top: float,
    width: float,
    height: float,
    font_name: str | None = None,
    layout_font_defaults: dict[str, Any] | None = None,
) -> Any:
    # ------------------------------------------------------------
    # 入力確認
    # ------------------------------------------------------------
    if table_definition.is_empty:
        raise ValueError(
            "表データが空です．"
        )

    if width <= 0:
        raise ValueError(
            "表の幅は0より大きくしてください．"
        )

    if height <= 0:
        raise ValueError(
            "表の高さは0より大きくしてください．"
        )

    row_count = table_definition.row_count
    column_count = table_definition.column_count

    style = get_table_style(
        table_definition.style_key,
        theme=theme,
    )

    # ------------------------------------------------------------
    # 表題領域
    # ------------------------------------------------------------
    caption_height = (
        DEFAULT_CAPTION_HEIGHT
        if table_definition.caption
        else 0.0
    )

    if table_definition.caption:
        _add_caption(
            slide,
            text=table_definition.caption,
            left=left,
            top=top,
            width=width,
            height=caption_height,
            theme=theme,
            font_name=font_name,
            font_size=resolve_font_size(role="table_caption", theme=theme, layout_defaults=layout_font_defaults),
        )

    table_top = (
        top
        + caption_height
    )

    table_height = max(
        MIN_TABLE_HEIGHT,
        height - caption_height,
    )

    # ===== DEBUG START =====
    print("----- TABLE BOUNDS -----")
    print(f"{top=}")
    print(f"{height=}")
    print(f"{caption_height=}")
    print(f"{table_top=}")
    print(f"{table_height=}")
    print(f"expected_table_bottom={table_top + table_height}")
    # ===== DEBUG END =====

    # ------------------------------------------------------------
    # PowerPoint表作成
    # ------------------------------------------------------------
    table_shape = slide.shapes.add_table(
        row_count,
        column_count,
        Inches(left),
        Inches(table_top),
        Inches(width),
        Inches(table_height),
    )

    # ===== DEBUG START =====
    print("----- TABLE SHAPE -----")
    print(f"shape_top={table_shape.top.inches}")
    print(f"shape_height={table_shape.height.inches}")
    print(
        "shape_bottom="
        f"{table_shape.top.inches + table_shape.height.inches}"
    )
    # ===== DEBUG END =====
    
    pptx_table = table_shape.table

    _set_uniform_column_widths(
        pptx_table,
        total_width=width,
        header_column_count=(
            table_definition.header.column_count
        ),
    )

    _set_uniform_row_heights(
        pptx_table,
        total_height=table_height,
    )

    # ------------------------------------------------------------
    # 通常セルへ文字・スタイルを設定
    #
    # 結合前に設定する
    # ------------------------------------------------------------
    for row_index in range(
        row_count,
    ):
        for column_index in range(
            column_count,
        ):
            definition_cell = _get_definition_cell(
                table_definition,
                row_index=row_index,
                column_index=column_index,
            )

            pptx_cell = pptx_table.cell(
                row_index,
                column_index,
            )

            _set_cell_margins(
                pptx_cell,
                style=style,
            )

            fill_color = _get_cell_fill_color(
                table_definition,
                row_index=row_index,
                column_index=column_index,
                style=style,
            )

            text_color = _get_cell_text_color(
                table_definition,
                row_index=row_index,
                column_index=column_index,
                style=style,
            )

            is_header = _is_header_cell(
                table_definition,
                row_index=row_index,
                column_index=column_index,
            )

            is_top_header = _is_top_header_cell(
                table_definition,
                row_index=row_index,
            )

            is_row_header = _is_row_header_cell(
                table_definition,
                column_index=column_index,
            )

            # --------------------------------------------------------
            # セル罫線
            #
            # python-pptxでは，セル背景を設定した後に罫線XMLを
            # 追加すると，PowerPointで罫線が表示されない場合がある．
            # 先に罫線XMLを設定し，その後で背景色を設定する．
            # --------------------------------------------------------
            _set_cell_borders(
                pptx_cell,
                row_index=row_index,
                column_index=column_index,
                row_count=row_count,
                column_count=column_count,
                style=style,
                border_color=theme.panel_line_color,
            )

            # --------------------------------------------------------
            # セル背景
            # --------------------------------------------------------
            _set_cell_fill(
                pptx_cell,
                fill_color,
            )

            text = (
                definition_cell.text
                if (
                    definition_cell is not None
                    and not definition_cell.merge_direction
                )
                else ""
            )

            _set_cell_text(
                pptx_cell,
                text=text,
                theme=theme,
                font_size=(
                    table_definition.font_size
                    if table_definition.font_size is not None
                    else resolve_font_size(
                        role=(
                            "table_header"
                            if is_top_header
                            else "table_body"
                        ),
                        theme=theme,
                        layout_defaults=layout_font_defaults,
                    )
                ),
                font_name=font_name,
                color=text_color,
                bold=(
                    style.header_bold
                    if (
                        is_top_header
                        or is_row_header
                    )
                    else False
                ),
                alignment=_get_cell_alignment(
                    table_definition,
                    row_index=row_index,
                    column_index=column_index,
                    text=text,
                ),
            )

    # ------------------------------------------------------------
    # セル結合
    # ------------------------------------------------------------
    _merge_table_cells(
        pptx_table,
        table_definition=table_definition,
    )

    # --------------------------------------------------------
    # 表の実際の下端座標を返す
    # --------------------------------------------------------
    return table_shape