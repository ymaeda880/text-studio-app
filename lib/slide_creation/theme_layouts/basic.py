# -*- coding: utf-8 -*-
# text_studio_app/lib/slide_creation/theme_layouts/basic.py
# ============================================================
# basic 標準レイアウト
#
# 機能：
# - 標準テーマ共通のページレイアウトを定義する
# - タイトル，見出し，小見出し，本文，最終ページを描画する
#
# 対象テーマ：
# - business_blue
# - simple_light
# - modern_dark
# - academic_green
#
# 方針：
# - 色とフォントはSlideThemeから取得する
# - ページの構成，位置，大きさはこのファイルで管理する
# - PowerPointへの基本部品描画はcomponents.pyへ任せる
# ============================================================

from __future__ import annotations

# ============================================================
# imports
# ============================================================
from pathlib import Path
from typing import Any

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

from lib.slide_creation.components import (
    add_bullet_text,
    add_footer,
    add_header,
    add_panel,
    add_slide_image,
    add_textbox,
    set_slide_background,
    split_body_lines,
)
from lib.slide_creation.models import (
    PresentationSettings,
    SlideDefinition,
    SlideTheme,
)
from lib.slide_creation.table.renderer import (
    render_table,
)

# ============================================================
# スライドサイズ
# ============================================================
_SLIDE_WIDTH = Inches(13.333)
_SLIDE_HEIGHT = Inches(7.5)


# ============================================================
# タイトルページ
# ============================================================
def render_title(
    slide: Any,
    *,
    slide_def: SlideDefinition,
    settings: PresentationSettings,
    theme: SlideTheme,
) -> None:
    style_key = slide_def.style_key

    if style_key == "full_color":
        set_slide_background(
            slide,
            color=theme.primary_color,
        )

        add_textbox(
            slide,
            left=1.0,
            top=2.1,
            width=11.3,
            height=1.4,
            text=slide_def.title,
            theme=theme,
            font_size=38,
            color=theme.title_text_color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

        add_textbox(
            slide,
            left=1.3,
            top=3.6,
            width=10.7,
            height=0.7,
            text=slide_def.subtitle,
            theme=theme,
            font_size=20,
            color=theme.title_text_color,
            alignment=PP_ALIGN.CENTER,
        )

    elif style_key == "centered":
        set_slide_background(
            slide,
            color=theme.background_color,
        )

        add_textbox(
            slide,
            left=1.0,
            top=2.15,
            width=11.3,
            height=1.4,
            text=slide_def.title,
            theme=theme,
            font_size=38,
            color=theme.body_text_color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

        add_textbox(
            slide,
            left=1.4,
            top=3.65,
            width=10.5,
            height=0.7,
            text=slide_def.subtitle,
            theme=theme,
            font_size=19,
            color=theme.sub_text_color,
            alignment=PP_ALIGN.CENTER,
        )

    elif style_key == "vertical_info":
        set_slide_background(
            slide,
            color=theme.background_color,
        )

        # ----------------------------------------------------
        # タイトル
        # ----------------------------------------------------
        add_textbox(
            slide,
            left=0.85,
            top=0.45,
            width=11.63,
            height=1.10,
            text=slide_def.title,
            theme=theme,
            font_size=38,
            color=theme.body_text_color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

        # ----------------------------------------------------
        # サブタイトル
        # ----------------------------------------------------
        add_textbox(
            slide,
            left=1.20,
            top=1.75,
            width=10.93,
            height=0.70,
            text=slide_def.subtitle,
            theme=theme,
            font_size=19,
            color=theme.sub_text_color,
            alignment=PP_ALIGN.CENTER,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

        presenter_name = (
            slide_def.presenter_name
            or settings.presenter_name
        )

        # ----------------------------------------------------
        # 会社名
        # ----------------------------------------------------
        if settings.company_name:
            add_textbox(
                slide,
                left=3.15,
                top=5.00,
                width=7.03,
                height=0.35,
                text=settings.company_name,
                theme=theme,
                font_size=11,
                color=theme.sub_text_color,
                alignment=PP_ALIGN.CENTER,
                vertical_anchor=MSO_ANCHOR.MIDDLE,
            )

        # ----------------------------------------------------
        # 発表者名
        # ----------------------------------------------------
        if presenter_name:
            add_textbox(
                slide,
                left=3.15,
                top=5.45,
                width=7.03,
                height=0.42,
                text=presenter_name,
                theme=theme,
                font_size=14,
                color=theme.body_text_color,
                bold=True,
                alignment=PP_ALIGN.CENTER,
                vertical_anchor=MSO_ANCHOR.MIDDLE,
            )

        # ----------------------------------------------------
        # 発表日
        # ----------------------------------------------------
        if settings.presentation_date:
            add_textbox(
                slide,
                left=3.15,
                top=5.95,
                width=7.03,
                height=0.35,
                text=settings.presentation_date,
                theme=theme,
                font_size=11,
                color=theme.sub_text_color,
                alignment=PP_ALIGN.CENTER,
                vertical_anchor=MSO_ANCHOR.MIDDLE,
            )

    elif style_key == "minimal":
        set_slide_background(
            slide,
            color=theme.background_color,
        )

        add_textbox(
            slide,
            left=1.0,
            top=2.35,
            width=11.0,
            height=1.3,
            text=slide_def.title,
            theme=theme,
            font_size=36,
            color=theme.body_text_color,
            bold=True,
        )

        add_textbox(
            slide,
            left=1.0,
            top=3.75,
            width=11.0,
            height=0.7,
            text=slide_def.subtitle,
            theme=theme,
            font_size=18,
            color=theme.sub_text_color,
        )

    else:
        set_slide_background(
            slide,
            color=theme.background_color,
        )

        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0,
            0,
            Inches(0.3),
            _SLIDE_HEIGHT,
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = theme.accent_color
        accent.line.fill.background()

        add_textbox(
            slide,
            left=1.0,
            top=2.15,
            width=11.2,
            height=1.35,
            text=slide_def.title,
            theme=theme,
            font_size=36,
            color=theme.body_text_color,
            bold=True,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

        add_textbox(
            slide,
            left=1.0,
            top=3.75,
            width=11.0,
            height=0.7,
            text=slide_def.subtitle,
            theme=theme,
            font_size=18,
            color=theme.sub_text_color,
        )


    # ========================================================
    # 下部詳細情報
    # ========================================================
    # vertical_infoは上の専用分岐内で縦配置済みのため，
    # 共通の横並び表示は行わない．
    if style_key != "vertical_info":
        detail_values = [
            value
            for value in (
                settings.company_name,
                (
                    slide_def.presenter_name
                    or settings.presenter_name
                ),
                settings.presentation_date,
            )
            if value
        ]

        add_textbox(
            slide,
            left=1.0,
            top=6.25,
            width=11.0,
            height=0.4,
            text="　｜　".join(detail_values),
            theme=theme,
            font_size=11,
            color=(
                theme.title_text_color
                if style_key == "full_color"
                else theme.sub_text_color
            ),
            alignment=PP_ALIGN.RIGHT,
        )


# ============================================================
# 見出しページ
# ============================================================
def render_section(
    slide: Any,
    *,
    slide_def: SlideDefinition,
    theme: SlideTheme,
) -> None:
    style_key = slide_def.style_key

    if style_key == "full_color":
        set_slide_background(
            slide,
            color=theme.primary_color,
        )

        add_textbox(
            slide,
            left=1.0,
            top=1.6,
            width=11.3,
            height=0.8,
            text=slide_def.section_number,
            theme=theme,
            font_size=26,
            color=theme.accent_color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        add_textbox(
            slide,
            left=1.0,
            top=2.5,
            width=11.3,
            height=1.4,
            text=slide_def.title,
            theme=theme,
            font_size=40,
            color=theme.title_text_color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

        add_textbox(
            slide,
            left=1.5,
            top=4.05,
            width=10.3,
            height=0.8,
            text=slide_def.subtitle,
            theme=theme,
            font_size=17,
            color=theme.title_text_color,
            alignment=PP_ALIGN.CENTER,
        )
        return

    set_slide_background(
        slide,
        color=theme.background_color,
    )

    if style_key == "large_number":
        add_textbox(
            slide,
            left=0.8,
            top=1.0,
            width=4.0,
            height=2.5,
            text=slide_def.section_number,
            theme=theme,
            font_size=80,
            color=theme.accent_color,
            bold=True,
        )

        add_textbox(
            slide,
            left=3.6,
            top=2.2,
            width=8.6,
            height=1.2,
            text=slide_def.title,
            theme=theme,
            font_size=38,
            color=theme.body_text_color,
            bold=True,
        )

        add_textbox(
            slide,
            left=3.65,
            top=3.55,
            width=8.2,
            height=0.8,
            text=slide_def.subtitle,
            theme=theme,
            font_size=17,
            color=theme.sub_text_color,
        )
        return

    if style_key == "centered":
        add_textbox(
            slide,
            left=1.0,
            top=1.8,
            width=11.3,
            height=0.6,
            text=slide_def.section_number,
            theme=theme,
            font_size=22,
            color=theme.accent_color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        add_textbox(
            slide,
            left=1.0,
            top=2.55,
            width=11.3,
            height=1.2,
            text=slide_def.title,
            theme=theme,
            font_size=40,
            color=theme.body_text_color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        add_textbox(
            slide,
            left=1.4,
            top=3.9,
            width=10.5,
            height=0.8,
            text=slide_def.subtitle,
            theme=theme,
            font_size=17,
            color=theme.sub_text_color,
            alignment=PP_ALIGN.CENTER,
        )
        return

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.85),
        Inches(1.25),
        Inches(0.12),
        Inches(4.8),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = theme.accent_color
    accent.line.fill.background()

    add_textbox(
        slide,
        left=1.3,
        top=1.6,
        width=10.5,
        height=0.6,
        text=slide_def.section_number,
        theme=theme,
        font_size=22,
        color=theme.accent_color,
        bold=True,
    )

    add_textbox(
        slide,
        left=1.3,
        top=2.35,
        width=10.5,
        height=1.2,
        text=slide_def.title,
        theme=theme,
        font_size=38,
        color=theme.body_text_color,
        bold=True,
    )

    add_textbox(
        slide,
        left=1.3,
        top=3.75,
        width=10.3,
        height=0.9,
        text=slide_def.subtitle,
        theme=theme,
        font_size=17,
        color=theme.sub_text_color,
    )


# ============================================================
# 小見出しページ
# ============================================================
def render_subsection(
    slide: Any,
    *,
    slide_def: SlideDefinition,
    theme: SlideTheme,
) -> None:
    style_key = slide_def.style_key

    set_slide_background(
        slide,
        color=theme.background_color,
    )

    if style_key == "accent_band":
        band = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0,
            Inches(2.2),
            _SLIDE_WIDTH,
            Inches(2.1),
        )
        band.fill.solid()
        band.fill.fore_color.rgb = theme.primary_color
        band.line.fill.background()

        add_textbox(
            slide,
            left=1.0,
            top=2.45,
            width=11.3,
            height=0.9,
            text=slide_def.title,
            theme=theme,
            font_size=34,
            color=theme.title_text_color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

        add_textbox(
            slide,
            left=1.5,
            top=3.35,
            width=10.3,
            height=0.6,
            text=slide_def.subtitle,
            theme=theme,
            font_size=16,
            color=theme.title_text_color,
            alignment=PP_ALIGN.CENTER,
        )
        return

    if style_key == "card":
        add_panel(
            slide,
            left=1.15,
            top=1.65,
            width=11.0,
            height=4.1,
            theme=theme,
        )

        add_textbox(
            slide,
            left=1.65,
            top=2.25,
            width=10.0,
            height=1.0,
            text=slide_def.title,
            theme=theme,
            font_size=34,
            color=theme.body_text_color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        add_textbox(
            slide,
            left=2.0,
            top=3.55,
            width=9.3,
            height=0.8,
            text=slide_def.subtitle,
            theme=theme,
            font_size=17,
            color=theme.sub_text_color,
            alignment=PP_ALIGN.CENTER,
        )
        return

    if style_key == "dark_panel":
        add_panel(
            slide,
            left=0.85,
            top=1.55,
            width=11.65,
            height=4.3,
            theme=theme,
            fill_color=theme.primary_color,
        )

        add_textbox(
            slide,
            left=1.35,
            top=2.25,
            width=10.7,
            height=1.0,
            text=slide_def.title,
            theme=theme,
            font_size=34,
            color=theme.title_text_color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        add_textbox(
            slide,
            left=1.7,
            top=3.55,
            width=10.0,
            height=0.8,
            text=slide_def.subtitle,
            theme=theme,
            font_size=17,
            color=theme.title_text_color,
            alignment=PP_ALIGN.CENTER,
        )
        return

    add_textbox(
        slide,
        left=1.0,
        top=2.4,
        width=11.3,
        height=1.0,
        text=slide_def.title,
        theme=theme,
        font_size=36,
        color=theme.body_text_color,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    add_textbox(
        slide,
        left=1.5,
        top=3.65,
        width=10.3,
        height=0.8,
        text=slide_def.subtitle,
        theme=theme,
        font_size=17,
        color=theme.sub_text_color,
        alignment=PP_ALIGN.CENTER,
    )


# ============================================================
# 本文ページ
# ============================================================
def render_content(
    slide: Any,
    *,
    slide_def: SlideDefinition,
    settings: PresentationSettings,
    page_number: int,
    theme: SlideTheme,
    inbox_root: Path,
    sub: str,
) -> None:
    set_slide_background(
        slide,
        color=theme.background_color,
    )

    add_header(
        slide,
        slide_width=_SLIDE_WIDTH,
        slide_title=slide_def.title,
        header_key=settings.header_key,
        theme=theme,
    )

    lines = split_body_lines(slide_def.body)

    # --------------------------------------------------------
    # 表
    #
    # 配置順：
    # - description（任意）
    # - table
    # - note（任意）
    # --------------------------------------------------------
    if slide_def.style_key == "table":
        if slide_def.table is None:
            raise ValueError(
                "style=tableですが，"
                "表データが設定されていません．"
            )

        # ----------------------------------------------------
        # 表全体のパネル
        # ----------------------------------------------------
        add_panel(
            slide,
            left=0.75,
            top=1.40,
            width=11.85,
            height=4.85,
            theme=theme,
        )

        # ----------------------------------------------------
        # パネル内部の配置領域
        # ----------------------------------------------------
        inner_left = 1.05
        inner_top = 1.62
        inner_width = 11.25
        inner_bottom = 6.02

        description_height = (
            0.62
            if slide_def.description
            else 0.0
        )

        note_height = (
            0.52
            if slide_def.note
            else 0.0
        )

        description_gap = (
            0.10
            if slide_def.description
            else 0.0
        )

        note_gap = (
            0.10
            if slide_def.note
            else 0.0
        )

        # ----------------------------------------------------
        # 説明文
        # ----------------------------------------------------
        if slide_def.description:
            add_textbox(
                slide,
                left=inner_left,
                top=inner_top,
                width=inner_width,
                height=description_height,
                text=slide_def.description,
                theme=theme,
                font_size=14,
                color=theme.body_text_color,
                vertical_anchor=MSO_ANCHOR.MIDDLE,
            )

        table_top = (
            inner_top
            + description_height
            + description_gap
        )

        table_bottom = (
            inner_bottom
            - note_height
            - note_gap
        )

        table_height = (
            table_bottom
            - table_top
        )

        # ----------------------------------------------------
        # 表
        # ----------------------------------------------------
        render_table(
            slide,
            table_definition=slide_def.table,
            theme=theme,
            left=inner_left,
            top=table_top,
            width=inner_width,
            height=table_height,
        )

        # ----------------------------------------------------
        # 補足
        # ----------------------------------------------------
        if slide_def.note:
            add_textbox(
                slide,
                left=inner_left,
                top=table_bottom + note_gap,
                width=inner_width,
                height=note_height,
                text=slide_def.note,
                theme=theme,
                font_size=12,
                color=theme.sub_text_color,
                vertical_anchor=MSO_ANCHOR.MIDDLE,
            )

    # --------------------------------------------------------
    # 左文章・右表
    #
    # 左側：
    # - itemizeまたは本文
    #
    # 右側：
    # - description（任意）
    # - table
    # - note（任意）
    # --------------------------------------------------------
    elif slide_def.style_key == "text_table":
        if slide_def.table is None:
            raise ValueError(
                "style=text_tableですが，"
                "表データが設定されていません．"
            )

        # ----------------------------------------------------
        # 左側文章パネル
        # ----------------------------------------------------
        add_panel(
            slide,
            left=0.75,
            top=1.40,
            width=5.85,
            height=4.85,
            theme=theme,
        )

        add_bullet_text(
            slide,
            lines=lines,
            left=1.05,
            top=1.75,
            width=5.20,
            height=4.00,
            theme=theme,
        )

        # ----------------------------------------------------
        # 右側表パネル
        # ----------------------------------------------------
        add_panel(
            slide,
            left=6.75,
            top=1.40,
            width=5.85,
            height=4.85,
            theme=theme,
        )

        # ----------------------------------------------------
        # 右側パネル内部の配置領域
        # ----------------------------------------------------
        table_area_left = 7.05
        table_area_top = 1.62
        table_area_width = 5.25
        table_area_bottom = 6.02

        description_height = (
            0.62
            if slide_def.description
            else 0.0
        )

        note_height = (
            0.52
            if slide_def.note
            else 0.0
        )

        description_gap = (
            0.10
            if slide_def.description
            else 0.0
        )

        note_gap = (
            0.10
            if slide_def.note
            else 0.0
        )

        # ----------------------------------------------------
        # 表の説明文
        # ----------------------------------------------------
        if slide_def.description:
            add_textbox(
                slide,
                left=table_area_left,
                top=table_area_top,
                width=table_area_width,
                height=description_height,
                text=slide_def.description,
                theme=theme,
                font_size=12,
                color=theme.body_text_color,
                vertical_anchor=MSO_ANCHOR.MIDDLE,
            )

        table_top = (
            table_area_top
            + description_height
            + description_gap
        )

        table_bottom = (
            table_area_bottom
            - note_height
            - note_gap
        )

        table_height = (
            table_bottom
            - table_top
        )

        # ----------------------------------------------------
        # 表
        # ----------------------------------------------------
        render_table(
            slide,
            table_definition=slide_def.table,
            theme=theme,
            left=table_area_left,
            top=table_top,
            width=table_area_width,
            height=table_height,
        )

        # ----------------------------------------------------
        # 補足
        # ----------------------------------------------------
        if slide_def.note:
            add_textbox(
                slide,
                left=table_area_left,
                top=table_bottom + note_gap,
                width=table_area_width,
                height=note_height,
                text=slide_def.note,
                theme=theme,
                font_size=11,
                color=theme.sub_text_color,
                vertical_anchor=MSO_ANCHOR.MIDDLE,
            )

    # --------------------------------------------------------
    # 左文章・右画像
    # --------------------------------------------------------
    elif slide_def.style_key == "text_image":

        # ----------------------------------------------------
        # 左側文章パネル
        # ----------------------------------------------------
        add_panel(
            slide,
            left=0.75,
            top=1.4,
            width=5.85,
            height=4.85,
            theme=theme,
        )

        add_bullet_text(
            slide,
            lines=lines,
            left=1.05,
            top=1.75,
            width=5.2,
            height=4.0,
            theme=theme,
        )

        # ----------------------------------------------------
        # 右側画像パネル
        # ----------------------------------------------------
        add_panel(
            slide,
            left=6.75,
            top=1.4,
            width=5.85,
            height=4.85,
            theme=theme,
        )

        # ----------------------------------------------------
        # 画像
        # ----------------------------------------------------
        image_height = (
            3.65
            if slide_def.image_caption
            else 4.15
        )

        add_slide_image(
            slide,
            image_file=slide_def.image_file,
            image_path=settings.image_path,
            inbox_root=inbox_root,
            sub=sub,
            left=7.05,
            top=1.70,
            width=5.25,
            height=image_height,
        )

        # ----------------------------------------------------
        # 画像説明
        # ----------------------------------------------------
        if slide_def.image_caption:
            add_textbox(
                slide,
                left=7.05,
                top=5.48,
                width=5.25,
                height=0.45,
                text=slide_def.image_caption,
                theme=theme,
                font_size=12,
                color=theme.sub_text_color,
                alignment=PP_ALIGN.CENTER,
                vertical_anchor=MSO_ANCHOR.MIDDLE,
            )

    elif slide_def.style_key == "two_column":   

        split_index = (len(lines) + 1) // 2
        left_lines = lines[:split_index]
        right_lines = lines[split_index:]

        add_panel(
            slide,
            left=0.75,
            top=1.4,
            width=5.85,
            height=4.85,
            theme=theme,
        )
        add_panel(
            slide,
            left=6.75,
            top=1.4,
            width=5.85,
            height=4.85,
            theme=theme,
        )

        add_bullet_text(
            slide,
            lines=left_lines,
            left=1.05,
            top=1.75,
            width=5.2,
            height=4.0,
            theme=theme,
        )
        add_bullet_text(
            slide,
            lines=right_lines,
            left=7.05,
            top=1.75,
            width=5.2,
            height=4.0,
            theme=theme,
        )

    elif slide_def.style_key == "comparison":
        left_lines = lines[0::2]
        right_lines = lines[1::2]

        add_textbox(
            slide,
            left=0.9,
            top=1.25,
            width=5.5,
            height=0.35,
            text=slide_def.left_heading or "項目A",
            theme=theme,
            font_size=16,
            color=theme.accent_color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        add_textbox(
            slide,
            left=6.9,
            top=1.25,
            width=5.5,
            height=0.35,
            text=slide_def.right_heading or "項目B",
            theme=theme,
            font_size=16,
            color=theme.accent_color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        add_panel(
            slide,
            left=0.75,
            top=1.65,
            width=5.85,
            height=4.6,
            theme=theme,
        )
        add_panel(
            slide,
            left=6.75,
            top=1.65,
            width=5.85,
            height=4.6,
            theme=theme,
        )

        add_bullet_text(
            slide,
            lines=left_lines,
            left=1.05,
            top=1.95,
            width=5.2,
            height=3.9,
            theme=theme,
        )
        add_bullet_text(
            slide,
            lines=right_lines,
            left=7.05,
            top=1.95,
            width=5.2,
            height=3.9,
            theme=theme,
        )

    elif slide_def.style_key == "message":
        add_panel(
            slide,
            left=1.0,
            top=1.65,
            width=11.3,
            height=4.4,
            theme=theme,
        )

        add_textbox(
            slide,
            left=1.5,
            top=2.2,
            width=10.3,
            height=3.0,
            text=slide_def.body,
            theme=theme,
            font_size=30,
            color=theme.body_text_color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

    else:
        add_panel(
            slide,
            left=0.8,
            top=1.4,
            width=11.7,
            height=4.85,
            theme=theme,
        )

        add_bullet_text(
            slide,
            lines=lines,
            left=1.15,
            top=1.75,
            width=10.9,
            height=4.1,
            theme=theme,
        )

    add_footer(
        slide,
        footer_key=settings.footer_key,
        settings=settings,
        page_number=page_number,
        theme=theme,
    )


# ============================================================
# 最終ページ
# ============================================================
def render_ending(
    slide: Any,
    *,
    slide_def: SlideDefinition,
    settings: PresentationSettings,
    theme: SlideTheme,
) -> None:
    style_key = slide_def.style_key

    if style_key == "thank_you":
        set_slide_background(
            slide,
            color=theme.primary_color,
        )

        add_textbox(
            slide,
            left=1.0,
            top=2.15,
            width=11.3,
            height=1.4,
            text=(
                slide_def.title
                or "ご清聴ありがとうございました"
            ),
            theme=theme,
            font_size=38,
            color=theme.title_text_color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

        add_textbox(
            slide,
            left=1.4,
            top=3.75,
            width=10.5,
            height=0.8,
            text=slide_def.subtitle,
            theme=theme,
            font_size=18,
            color=theme.title_text_color,
            alignment=PP_ALIGN.CENTER,
        )
        return

    set_slide_background(
        slide,
        color=theme.background_color,
    )

    if style_key == "contact":
        add_textbox(
            slide,
            left=1.0,
            top=1.25,
            width=11.3,
            height=0.8,
            text=slide_def.title or "お問い合わせ",
            theme=theme,
            font_size=34,
            color=theme.body_text_color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        add_panel(
            slide,
            left=1.4,
            top=2.35,
            width=10.5,
            height=3.2,
            theme=theme,
        )

        contact_values = [
            value
            for value in (
                settings.company_name,
                (
                    slide_def.presenter_name
                    or settings.presenter_name
                ),
                slide_def.contact_text,
            )
            if value
        ]

        add_textbox(
            slide,
            left=1.9,
            top=2.75,
            width=9.5,
            height=2.3,
            text="\n".join(contact_values),
            theme=theme,
            font_size=20,
            color=theme.body_text_color,
            alignment=PP_ALIGN.CENTER,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )
        return

    if style_key == "summary":
        add_textbox(
            slide,
            left=0.75,
            top=0.55,
            width=11.8,
            height=0.7,
            text=slide_def.title or "まとめ",
            theme=theme,
            font_size=32,
            color=theme.body_text_color,
            bold=True,
        )

        add_panel(
            slide,
            left=0.85,
            top=1.55,
            width=11.6,
            height=4.9,
            theme=theme,
        )

        add_bullet_text(
            slide,
            lines=split_body_lines(slide_def.body),
            left=1.25,
            top=1.95,
            width=10.8,
            height=4.0,
            theme=theme,
        )
        return

    add_textbox(
        slide,
        left=1.0,
        top=2.45,
        width=11.3,
        height=1.2,
        text=slide_def.title or "以上です",
        theme=theme,
        font_size=38,
        color=theme.body_text_color,
        bold=True,
        alignment=PP_ALIGN.CENTER,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )

    add_textbox(
        slide,
        left=1.4,
        top=3.8,
        width=10.5,
        height=0.8,
        text=slide_def.subtitle,
        theme=theme,
        font_size=18,
        color=theme.sub_text_color,
        alignment=PP_ALIGN.CENTER,
    )