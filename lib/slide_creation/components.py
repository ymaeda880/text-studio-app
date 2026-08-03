# -*- coding: utf-8 -*-
# text_studio_app/lib/slide_creation/components.py
# ============================================================
# PowerPointスライド 共通描画部品
#
# 機能：
# - 背景を設定する
# - テキストボックスを追加する
# - パネル，ヘッダー，フッターを追加する
# - 箇条書きテキストを追加する
# ============================================================

from __future__ import annotations

# ============================================================
# imports
# ============================================================
from pathlib import Path
from typing import Any

from PIL import Image, ImageOps
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from lib.slide_creation.models import (
    PresentationSettings,
    SlideTheme,
)
from lib.slide_creation.slide_image_resolver import (
    resolve_slide_image_path,
)


# ============================================================
# 背景
# ============================================================
def set_slide_background(
    slide: Any,
    *,
    color: Any,
) -> None:
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = color


# ============================================================
# テキストボックス
# ============================================================
def add_textbox(
    slide: Any,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    theme: SlideTheme,
    font_size: int,
    color: Any,
    bold: bool = False,
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
    vertical_anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
) -> Any:
    box = slide.shapes.add_textbox(
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )

    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = vertical_anchor

    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = alignment

    for run in paragraph.runs:
        run.font.name = theme.font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color

    return box


# ============================================================
# パネル
# ============================================================
def add_panel(
    slide: Any,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    theme: SlideTheme,
    rounded: bool = True,
    fill_color: Any | None = None,
) -> Any:
    shape_type = (
        MSO_SHAPE.ROUNDED_RECTANGLE
        if rounded
        else MSO_SHAPE.RECTANGLE
    )

    panel = slide.shapes.add_shape(
        shape_type,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )

    panel.fill.solid()
    panel.fill.fore_color.rgb = fill_color or theme.panel_color
    panel.line.color.rgb = theme.panel_line_color
    panel.line.width = Pt(1.2)

    return panel

# ============================================================
# 画像
# ============================================================
def add_slide_image(
    slide: Any,
    *,
    image_file: str,
    image_path: str,
    inbox_root: Path,
    sub: str,
    left: float,
    top: float,
    width: float,
    height: float,
) -> Any:
    """
    SlideTexで指定された画像をPowerPointへ配置する．

    画像パスは，次の優先順位で解決する．

    1．image_fileが絶対パス
    2．image_pathが空またはinboxの場合はinbox検索
    3．image_pathで指定されたフォルダーを基準に検索

    画像の縦横比は維持し，
    指定された領域内へ中央配置する．
    """

    # --------------------------------------------------------
    # 画像パス解決
    # --------------------------------------------------------
    resolved_path = resolve_slide_image_path(
        image_file=image_file,
        image_path=image_path,
        inbox_root=inbox_root,
        sub=sub,
    )

    # --------------------------------------------------------
    # 画像サイズ取得
    # --------------------------------------------------------
    try:
        with Image.open(resolved_path) as source_image:
            corrected_image = ImageOps.exif_transpose(
                source_image
            )
            image_width_px, image_height_px = (
                corrected_image.size
            )

    except Exception as exc:
        raise ValueError(
            "画像サイズを取得できません："
            f"{resolved_path} / {exc}"
        ) from exc

    if image_width_px <= 0 or image_height_px <= 0:
        raise ValueError(
            "画像サイズが不正です："
            f"{resolved_path}"
        )

    # --------------------------------------------------------
    # 配置領域内へ縦横比を維持して収める
    # --------------------------------------------------------
    image_ratio = (
        float(image_width_px)
        / float(image_height_px)
    )
    box_ratio = float(width) / float(height)

    if image_ratio >= box_ratio:
        placed_width = float(width)
        placed_height = placed_width / image_ratio

    else:
        placed_height = float(height)
        placed_width = placed_height * image_ratio

    # --------------------------------------------------------
    # 配置領域の中央へ位置調整
    # --------------------------------------------------------
    placed_left = (
        float(left)
        + (
            float(width)
            - placed_width
        )
        / 2.0
    )
    placed_top = (
        float(top)
        + (
            float(height)
            - placed_height
        )
        / 2.0
    )

    # --------------------------------------------------------
    # PowerPointへ画像を追加
    # --------------------------------------------------------
    try:
        return slide.shapes.add_picture(
            str(resolved_path),
            Inches(placed_left),
            Inches(placed_top),
            width=Inches(placed_width),
            height=Inches(placed_height),
        )

    except Exception as exc:
        raise RuntimeError(
            "画像をPowerPointへ貼り込めません："
            f"{resolved_path} / {exc}"
        ) from exc
    
# ============================================================
# 箇条書き
# ============================================================
def add_bullet_text(
    slide: Any,
    *,
    lines: list[str],
    left: float,
    top: float,
    width: float,
    height: float,
    theme: SlideTheme,
    font_size: int | None = None,
) -> None:
    box = slide.shapes.add_textbox(
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )

    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP

    display_lines = lines or ["本文が入力されていません．"]

    for index, line in enumerate(display_lines):
        paragraph = (
            frame.paragraphs[0]
            if index == 0
            else frame.add_paragraph()
        )

        paragraph.text = f"●  {line}"
        paragraph.space_after = Pt(12)

        for run in paragraph.runs:
            run.font.name = theme.font_name
            run.font.size = Pt(font_size or theme.body_font_size)
            run.font.color.rgb = theme.body_text_color


# ============================================================
# 本文ページ用ヘッダー
# ============================================================
def add_header(
    slide: Any,
    *,
    slide_width: Any,
    slide_title: str,
    header_key: str,
    theme: SlideTheme,
) -> None:
    if header_key == "title_band":
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0,
            0,
            slide_width,
            Inches(1.05),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = theme.primary_color
        bar.line.fill.background()

        add_textbox(
            slide,
            left=0.65,
            top=0.17,
            width=11.8,
            height=0.7,
            text=slide_title,
            theme=theme,
            font_size=theme.title_font_size,
            color=theme.title_text_color,
            bold=True,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )
        return

    add_textbox(
        slide,
        left=0.65,
        top=0.28,
        width=11.8,
        height=0.7,
        text=slide_title,
        theme=theme,
        font_size=theme.title_font_size,
        color=theme.body_text_color,
        bold=True,
    )

    if header_key == "simple_line":
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.65),
            Inches(1.02),
            Inches(2.1),
            Inches(0.07),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = theme.accent_color
        line.line.fill.background()


# ============================================================
# 本文ページ用フッター
# ============================================================
def add_footer(
    slide: Any,
    *,
    footer_key: str,
    settings: PresentationSettings,
    page_number: int,
    theme: SlideTheme,
) -> None:
    if footer_key == "none":
        return

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.65),
        Inches(6.62),
        Inches(12.0),
        Inches(0.015),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = theme.panel_line_color
    line.line.fill.background()

    if footer_key == "detailed":
        values = [
            value
            for value in (
                settings.company_name,
                settings.presentation_title,
                settings.presentation_date,
            )
            if value
        ]
        left_text = "　｜　".join(values)

    elif footer_key == "confidential":
        left_text = (
            f"{settings.presentation_title}　｜　社外秘"
        )

    else:
        left_text = settings.presentation_title

    add_textbox(
        slide,
        left=0.7,
        top=6.68,
        width=10.7,
        height=0.3,
        text=left_text,
        theme=theme,
        font_size=theme.footer_font_size,
        color=theme.sub_text_color,
    )

    add_textbox(
        slide,
        left=11.7,
        top=6.68,
        width=0.6,
        height=0.3,
        text=str(page_number),
        theme=theme,
        font_size=theme.footer_font_size,
        color=theme.sub_text_color,
        alignment=PP_ALIGN.RIGHT,
    )


# ============================================================
# 文字列分割
# ============================================================
def split_body_lines(body: str) -> list[str]:
    return [
        line.strip()
        for line in body.splitlines()
        if line.strip()
    ]