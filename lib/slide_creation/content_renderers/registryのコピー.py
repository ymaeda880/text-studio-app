# -*- coding: utf-8 -*-
# text_studio_app/lib/slide_creation/content_renderers/registry.py
# ============================================================
# 本文コンテンツ 共通描画
# ============================================================
from __future__ import annotations
from pathlib import Path
from typing import Any, Mapping
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from lib.slide_creation.components import add_bullet_text, add_slide_image, add_textbox
from lib.slide_creation.content_layouts.models import RegionBounds
from lib.slide_creation.font_settings import resolve_font_name, resolve_font_size
from lib.slide_creation.models import ContentDefinition, PresentationSettings, RegionDefinition, SlideDefinition, SlideTheme
from lib.slide_creation.table.renderer import render_table

_INNER_MARGIN = 0.30
_CONTENT_GAP = 0.14

# ============================================================
# 表とnoteの間隔
# ============================================================
_TABLE_NOTE_GAP = 0.18

def _inner_bounds(bounds: RegionBounds)->RegionBounds:
    return RegionBounds(bounds.left+_INNER_MARGIN,bounds.top+0.22,max(bounds.width-_INNER_MARGIN*2,0.20),max(bounds.height-0.44,0.20))

def _render_text(slide:Any,*,content:ContentDefinition,bounds:RegionBounds,theme:SlideTheme,
                 slide_def:SlideDefinition,region:RegionDefinition,layout_defaults:Mapping[str,Any])->None:
    lines=[line.strip() for line in content.text.splitlines() if line.strip()]
    role="content_narrow_body" if bounds.width < 4.2 else "content_body"
    add_bullet_text(slide,lines=lines,left=bounds.left,top=bounds.top,width=bounds.width,height=bounds.height,
        theme=theme,font_size=resolve_font_size(role=role,theme=theme,slide_def=slide_def,region=region,layout_defaults=layout_defaults),
        font_name=resolve_font_name(theme=theme,slide_def=slide_def,region=region,layout_defaults=layout_defaults))

def _render_image(slide:Any,*,content:ContentDefinition,bounds:RegionBounds,settings:PresentationSettings,
                  inbox_root:Path,sub:str,theme:SlideTheme,slide_def:SlideDefinition,region:RegionDefinition,
                  layout_defaults:Mapping[str,Any])->None:
    caption_height=0.45 if content.image_caption else 0.0
    caption_gap=0.08 if content.image_caption else 0.0
    add_slide_image(slide,image_file=content.image_file,image_path=settings.image_path,inbox_root=inbox_root,sub=sub,
        left=bounds.left,top=bounds.top,width=bounds.width,height=max(bounds.height-caption_height-caption_gap,0.40))
    if content.image_caption:
        add_textbox(slide,left=bounds.left,top=bounds.top+bounds.height-caption_height,width=bounds.width,height=caption_height,
            text=content.image_caption,theme=theme,
            font_size=resolve_font_size(role="image_caption",theme=theme,slide_def=slide_def,region=region,layout_defaults=layout_defaults),
            font_name=resolve_font_name(theme=theme,slide_def=slide_def,region=region,layout_defaults=layout_defaults),
            color=theme.sub_text_color,alignment=PP_ALIGN.CENTER,vertical_anchor=MSO_ANCHOR.MIDDLE)

def _render_table_content(
    slide: Any,
    *,
    content: ContentDefinition,
    bounds: RegionBounds,
    theme: SlideTheme,
    slide_def: SlideDefinition,
    region: RegionDefinition,
    layout_defaults: Mapping[str, Any],
) -> None:
    if content.table is None:
        raise ValueError(
            "tableコンテンツですが，表データがありません．"
        )

    # --------------------------------------------------------
    # 配置領域
    # --------------------------------------------------------
    description_height = (
        0.52
        if content.description
        else 0.0
    )

    note_height = (
        0.44
        if content.note
        else 0.0
    )

    description_gap = (
        0.08
        if content.description
        else 0.0
    )

    note_gap = (
        _TABLE_NOTE_GAP
        if content.note
        else 0.0
    )
    # ===== DEBUG START =====
    print(f"{note_gap=}")
    print(f"{description_height=}")
    print(f"{note_height=}")
    # ===== DEBUG END =====

    # --------------------------------------------------------
    # フォント
    # --------------------------------------------------------
    font_name = resolve_font_name(
        theme=theme,
        slide_def=slide_def,
        region=region,
        layout_defaults=layout_defaults,
    )

    # --------------------------------------------------------
    # 表の説明文
    # --------------------------------------------------------
    if content.description:
        add_textbox(
            slide,
            left=bounds.left,
            top=bounds.top,
            width=bounds.width,
            height=description_height,
            text=content.description,
            theme=theme,
            font_size=resolve_font_size(
                role="table_description",
                theme=theme,
                slide_def=slide_def,
                region=region,
                layout_defaults=layout_defaults,
            ),
            font_name=font_name,
            color=theme.body_text_color,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

    # --------------------------------------------------------
    # 表の配置領域
    # --------------------------------------------------------
    table_top = (
        bounds.top
        + description_height
        + description_gap
    )

    table_height = max(
        bounds.height
        - description_height
        - description_gap
        - note_height
        - note_gap,
        0.60,
    )

    # --------------------------------------------------------
    # 表
    # --------------------------------------------------------
    table_shape = render_table(
        slide,
        table_definition=content.table,
        theme=theme,
        left=bounds.left,
        top=table_top,
        width=bounds.width,
        height=table_height,
        font_name=font_name,
        layout_font_defaults=layout_defaults,
    )

    # --------------------------------------------------------
    # 表の補足
    #
    # Region下端ではなく，
    # 実際に作成された表の下端を基準に配置する．
    # --------------------------------------------------------
    # if content.note:
    #     actual_table_bottom = (
    #         table_shape.top.inches
    #         + table_shape.height.inches
    #     )

    #     note_top = (
    #         actual_table_bottom
    #         + _TABLE_NOTE_GAP
    #     )

    #     # ===== DEBUG START =====
    #     print(f"{actual_table_bottom=}")
    #     print(f"{note_top=}")
    #     print(f"{_TABLE_NOTE_GAP=}")
    #     # ====

    #     # ===== DEBUG START =====
    #     print("note textbox")
    #     print(f"top={note_top}")
    #     print(f"height={note_height}")
    #     # ===== DEBUG END =====

    #     add_textbox(
    #         slide,
    #         left=bounds.left,
    #         top=note_top,
    #         width=bounds.width,
    #         height=note_height,
    #         text=content.note,
    #         theme=theme,
    #         font_size=resolve_font_size(
    #             role="table_note",
    #             theme=theme,
    #             slide_def=slide_def,
    #             region=region,
    #             layout_defaults=layout_defaults,
    #         ),
    #         font_name=font_name,
    #         color=theme.sub_text_color,
    #         vertical_anchor=MSO_ANCHOR.TOP,
    #     )
    # --------------------------------------------------------
    # 表の補足
    #
    # 表図形の下端を取得し，
    # その直下へnote文字を配置する．
    # --------------------------------------------------------
    if content.note:
        actual_table_bottom = (
            table_shape.top.inches
            + table_shape.height.inches
        )

        note_top = (
            actual_table_bottom
            + _TABLE_NOTE_GAP
        )

        note_box = slide.shapes.add_textbox(
            Inches(bounds.left),
            Inches(note_top),
            Inches(bounds.width),
            Inches(note_height),
        )

        note_frame = note_box.text_frame
        note_frame.clear()
        note_frame.word_wrap = True
        note_frame.vertical_anchor = MSO_ANCHOR.TOP

        # テキストボックス内部の余白をなくす
        note_frame.margin_left = Inches(0)
        note_frame.margin_right = Inches(0)
        note_frame.margin_top = Inches(0)
        note_frame.margin_bottom = Inches(0)

        paragraph = note_frame.paragraphs[0]
        paragraph.text = content.note
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)

        note_font_size = resolve_font_size(
            role="table_note",
            theme=theme,
            slide_def=slide_def,
            region=region,
            layout_defaults=layout_defaults,
        )

        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(note_font_size)
            run.font.color.rgb = theme.sub_text_color

def render_region_contents(slide:Any,*,contents:list[ContentDefinition],bounds:RegionBounds,settings:PresentationSettings,
                           inbox_root:Path,sub:str,theme:SlideTheme,slide_def:SlideDefinition,region:RegionDefinition,
                           layout_defaults:Mapping[str,Any])->None:
    if not contents:return
    inner=_inner_bounds(bounds); count=len(contents); total_gap=_CONTENT_GAP*max(count-1,0)
    content_height=max((inner.height-total_gap)/count,0.40)
    for index,content in enumerate(contents):
        cb=RegionBounds(inner.left,inner.top+index*(content_height+_CONTENT_GAP),inner.width,content_height)
        if content.content_type=="text": _render_text(slide,content=content,bounds=cb,theme=theme,slide_def=slide_def,region=region,layout_defaults=layout_defaults)
        elif content.content_type=="image": _render_image(slide,content=content,bounds=cb,settings=settings,inbox_root=inbox_root,sub=sub,theme=theme,slide_def=slide_def,region=region,layout_defaults=layout_defaults)
        elif content.content_type=="table": _render_table_content(slide,content=content,bounds=cb,theme=theme,slide_def=slide_def,region=region,layout_defaults=layout_defaults)
        else: raise ValueError("未対応の本文コンテンツです："+content.content_type)
